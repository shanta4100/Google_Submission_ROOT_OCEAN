#!/usr/bin/env python3
"""
generate_reports.py
-------------------
Master automation script for GNAIAAAC Root Ocean / Pure Water project.
Generates all key deliverables:
- PDF report
- DOCX master document
- Excel budget sheet
- PPTX executive deck
"""

import os
from pathlib import Path

# --- PDF generation ---
from fpdf import FPDF  # pip install fpdf

# --- DOCX generation ---
from docx import Document  # pip install python-docx

# --- Excel generation ---
import openpyxl  # pip install openpyxl

# --- PPTX generation ---
from pptx import Presentation  # pip install python-pptx

# --- Define output paths ---
OUTPUT_DIR = Path("output")
OUTPUT_DIR.mkdir(exist_ok=True)

PDF_FILE = OUTPUT_DIR / "GNAIAAAC_MASTER_ROOT.pdf"
DOCX_FILE = OUTPUT_DIR / "ROOT_OCEAN_MASTER.docx"
EXCEL_FILE = OUTPUT_DIR / "GNAIAAAC_Budget.xlsx"
PPTX_FILE = OUTPUT_DIR / "GNAIAAAC_Executive_Deck.pptx"

# ---------------- PDF REPORT ----------------
def generate_pdf():
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", "B", 16)
    pdf.cell(0, 10, "GNAIAAAC ROOT OCEAN / PURE WATER MASTER REPORT", ln=True, align="C")
    pdf.ln(10)
    pdf.set_font("Arial", "", 12)
    pdf.multi_cell(0, 8, "This report contains the full technical documentation, simulations, and analysis for the AI-driven water governance project.")
    pdf.output(str(PDF_FILE))
    print(f"PDF generated at {PDF_FILE}")

# ---------------- DOCX MASTER ----------------
def generate_docx():
    doc = Document()
    doc.add_heading("ROOT OCEAN MASTER DOCUMENT", level=0)
    doc.add_paragraph("This master document contains all text, references, and technical notes for the GNAIAAAC Root Ocean / Pure Water project.")
    doc.save(DOCX_FILE)
    print(f"DOCX generated at {DOCX_FILE}")

# ---------------- EXCEL BUDGET ----------------
def generate_excel():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Budget Plan"
    ws.append(["Item", "Description", "Estimated Cost (USD)"])
    ws.append(["Water Sensors", "AI monitoring sensors", 5000])
    ws.append(["Software Licenses", "AI and visualization tools", 2000])
    ws.append(["Staffing / Consulting", "Technical experts (remote)", 3000])
    wb.save(EXCEL_FILE)
    print(f"Excel generated at {EXCEL_FILE}")

# ---------------- PPTX DECK ----------------
def generate_pptx():
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "GNAIAAAC Root Ocean / Pure Water"
    slide.placeholders[1].text = "Executive Pitch Deck"
    
    # Add a content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Project Overview"
    slide.placeholders[1].text = "AI-driven simulations, visualizations, and automated workflows for water governance."
    
    prs.save(PPTX_FILE)
    print(f"PPTX generated at {PPTX_FILE}")

# ---------------- MAIN ----------------
if __name__ == "__main__":
    generate_pdf()
    generate_docx()
    generate_excel()
    generate_pptx()
    print("\nAll files generated successfully in the 'output' folder.")